import io
import json
import os

import anthropic
import streamlit as st
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# -----------------------------
# Page config
# -----------------------------
st.set_page_config(page_title="Deck Builder", layout="wide", page_icon="📊")

# -----------------------------
# Custom Styling
# -----------------------------
st.markdown("""
    <style>
    .stButton>button[kind="primary"] {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        border: none;
        border-radius: 8px;
        height: 3em;
        font-weight: 600;
    }
    .stButton>button { border-radius: 8px; height: 3em; font-weight: 500; }
    h1 {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        -webkit-background-clip: text; -webkit-text-fill-color: transparent;
    }
    </style>
""", unsafe_allow_html=True)

# -----------------------------
# Constants
# -----------------------------
SNAPLOGIC_PURPLE = RGBColor(0x4B, 0x2E, 0x83)
SNAPLOGIC_DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT_BLUE = RGBColor(0x66, 0x7E, 0xEA)

SNAPLOGIC_SOLUTIONS = [
    "Agent Creator (AI/LLM Agents)",
    "Data Integration & ETL/ELT",
    "API Management & Integration",
    "Data Pipeline Automation",
    "Cloud Data Migration",
    "Application Integration",
    "B2B / EDI Integration",
    "Self-Service Data Preparation",
    "Real-Time Data Streaming",
    "Intelligent Document Processing",
]

INDUSTRIES = [
    "Financial Services & Banking",
    "Healthcare & Life Sciences",
    "Retail & E-Commerce",
    "Manufacturing & Supply Chain",
    "Technology & SaaS",
    "Media & Entertainment",
    "Energy & Utilities",
    "Government & Public Sector",
    "Education",
    "Professional Services",
    "Other",
]


# -----------------------------
# Claude content generation
# -----------------------------
def generate_slide_content(
    customer_name: str,
    industry: str,
    pain_points: str,
    solutions: list[str],
    additional_context: str,
) -> dict:
    """Call Claude to generate structured slide content."""
    api_key = os.getenv("ANTHROPIC_API_KEY", "")
    client = anthropic.Anthropic(api_key=api_key)

    solutions_str = "\n".join(f"- {s}" for s in solutions)

    prompt = f"""You are a SnapLogic solutions expert and sales engineer creating a customer-facing PowerPoint presentation.

Customer: {customer_name}
Industry: {industry}
Pain Points / Challenges:
{pain_points}

Relevant SnapLogic Solutions:
{solutions_str}

Additional Context:
{additional_context if additional_context else "None provided."}

Generate content for a polished, customer-ready presentation with the following slides. Return ONLY valid JSON matching this exact structure:

{{
  "title_slide": {{
    "title": "...",
    "subtitle": "..."
  }},
  "agenda_slide": {{
    "title": "Today's Agenda",
    "items": ["...", "...", "...", "...", "..."]
  }},
  "customer_challenges_slide": {{
    "title": "...",
    "challenges": [
      {{"heading": "...", "detail": "..."}},
      {{"heading": "...", "detail": "..."}},
      {{"heading": "...", "detail": "..."}}
    ]
  }},
  "snaplogic_overview_slide": {{
    "title": "Why SnapLogic",
    "tagline": "...",
    "points": [
      {{"heading": "...", "detail": "..."}},
      {{"heading": "...", "detail": "..."}},
      {{"heading": "...", "detail": "..."}}
    ]
  }},
  "solution_slides": [
    {{
      "solution_name": "...",
      "title": "...",
      "value_prop": "...",
      "capabilities": ["...", "...", "..."],
      "customer_outcome": "..."
    }}
  ],
  "roi_slide": {{
    "title": "Expected Business Impact",
    "metrics": [
      {{"stat": "...", "description": "..."}},
      {{"stat": "...", "description": "..."}},
      {{"stat": "...", "description": "..."}}
    ],
    "summary": "..."
  }},
  "next_steps_slide": {{
    "title": "Recommended Next Steps",
    "steps": [
      {{"step": "Step 1", "action": "...", "timeline": "..."}},
      {{"step": "Step 2", "action": "...", "timeline": "..."}},
      {{"step": "Step 3", "action": "...", "timeline": "..."}}
    ]
  }},
  "closing_slide": {{
    "title": "...",
    "message": "...",
    "contact_prompt": "Let's connect and get started."
  }}
}}

Guidelines:
- Be specific to {customer_name} and their {industry} context
- Connect pain points directly to SnapLogic capabilities
- Use confident, customer-facing language
- Keep bullet points concise (under 12 words each)
- ROI stats should be realistic and industry-appropriate
- Create one solution_slide per selected solution (max 4 slides if more than 4 solutions selected)
"""

    message = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=4096,
        messages=[{"role": "user", "content": prompt}],
    )

    raw = message.content[0].text.strip()
    # Strip markdown code fences if present
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    return json.loads(raw)


# -----------------------------
# PowerPoint builder
# -----------------------------
def _set_slide_background(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, text, left, top, width, height, font_size, bold=False,
                  color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def _add_bullet_paragraph(tf, text, font_size=14, bold=False, color=WHITE, indent_level=1):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.level = indent_level
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def build_presentation(content: dict, customer_name: str) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # completely blank

    W = prs.slide_width
    H = prs.slide_height

    # ── Helpers ──────────────────────────────────────────────────────────────

    def add_slide():
        return prs.slides.add_slide(blank_layout)

    def purple_bg(slide):
        _set_slide_background(slide, SNAPLOGIC_PURPLE)

    def light_bg(slide):
        _set_slide_background(slide, LIGHT_GRAY)

    def add_rect(slide, left, top, width, height, fill_color, line_color=None):
        from pptx.util import Emu
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            left, top, width, height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if line_color:
            shape.line.color.rgb = line_color
        else:
            shape.line.fill.background()
        return shape

    def header_bar(slide, text, bg=SNAPLOGIC_DARK, fg=WHITE, font_size=28):
        add_rect(slide, 0, 0, W, Inches(0.85), bg)
        _add_text_box(slide, text,
                      Inches(0.4), Inches(0.1),
                      W - Inches(0.8), Inches(0.75),
                      font_size, bold=True, color=fg, align=PP_ALIGN.LEFT)

    def footer(slide, label="SnapLogic Confidential"):
        add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), SNAPLOGIC_DARK)
        _add_text_box(slide, label,
                      Inches(0.4), H - Inches(0.33),
                      W - Inches(0.8), Inches(0.3),
                      9, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.LEFT)
        _add_text_box(slide, "snaplogic.com",
                      0, H - Inches(0.33),
                      W - Inches(0.4), Inches(0.3),
                      9, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.RIGHT)

    # ── 1. Title Slide ───────────────────────────────────────────────────────
    slide = add_slide()
    purple_bg(slide)
    # Accent bar on left
    add_rect(slide, 0, 0, Inches(0.18), H, ACCENT_BLUE)
    # Logo area placeholder
    _add_text_box(slide, "SNAPLOGIC",
                  Inches(0.4), Inches(0.3),
                  Inches(4), Inches(0.5),
                  16, bold=True, color=WHITE)
    # Title
    d = content["title_slide"]
    _add_text_box(slide, d["title"],
                  Inches(0.4), Inches(1.8),
                  W - Inches(1.0), Inches(1.8),
                  38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    _add_text_box(slide, d["subtitle"],
                  Inches(0.4), Inches(3.8),
                  W - Inches(1.0), Inches(0.9),
                  20, color=RGBColor(0xCC, 0xCC, 0xFF), align=PP_ALIGN.LEFT)
    _add_text_box(slide, "Prepared exclusively for " + customer_name,
                  Inches(0.4), Inches(4.9),
                  W - Inches(1.0), Inches(0.5),
                  13, color=RGBColor(0xBB, 0xBB, 0xDD))
    footer(slide, "Confidential — Prepared for " + customer_name)

    # ── 2. Agenda Slide ──────────────────────────────────────────────────────
    slide = add_slide()
    light_bg(slide)
    header_bar(slide, content["agenda_slide"]["title"])
    footer(slide)
    add_rect(slide, Inches(0.4), Inches(1.0), Inches(0.08), H - Inches(1.5), ACCENT_BLUE)

    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.1), W - Inches(1.2), H - Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in content["agenda_slide"]["items"]:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"▸  {item}"
        run.font.size = Pt(18)
        run.font.color.rgb = SNAPLOGIC_DARK

    # ── 3. Customer Challenges Slide ─────────────────────────────────────────
    slide = add_slide()
    light_bg(slide)
    d = content["customer_challenges_slide"]
    header_bar(slide, d["title"])
    footer(slide)

    cols = len(d["challenges"])
    col_w = (W - Inches(0.8)) / cols
    for i, ch in enumerate(d["challenges"]):
        x = Inches(0.4) + i * col_w
        card = add_rect(slide, x + Inches(0.1), Inches(1.1),
                        col_w - Inches(0.2), H - Inches(1.9), SNAPLOGIC_PURPLE)
        _add_text_box(slide, ch["heading"],
                      x + Inches(0.25), Inches(1.25),
                      col_w - Inches(0.5), Inches(0.7),
                      15, bold=True, color=WHITE)
        _add_text_box(slide, ch["detail"],
                      x + Inches(0.25), Inches(2.1),
                      col_w - Inches(0.5), H - Inches(3.2),
                      13, color=RGBColor(0xDD, 0xDD, 0xFF))

    # ── 4. SnapLogic Overview Slide ──────────────────────────────────────────
    slide = add_slide()
    purple_bg(slide)
    d = content["snaplogic_overview_slide"]
    header_bar(slide, d["title"], bg=SNAPLOGIC_DARK)
    footer(slide)

    _add_text_box(slide, d["tagline"],
                  Inches(0.4), Inches(1.0),
                  W - Inches(0.8), Inches(0.6),
                  16, color=RGBColor(0xCC, 0xCC, 0xFF), align=PP_ALIGN.LEFT)

    for i, pt in enumerate(d["points"]):
        y = Inches(1.9) + i * Inches(1.5)
        add_rect(slide, Inches(0.4), y, Inches(0.06), Inches(1.1), ACCENT_BLUE)
        _add_text_box(slide, pt["heading"],
                      Inches(0.65), y,
                      W - Inches(1.1), Inches(0.45),
                      15, bold=True, color=WHITE)
        _add_text_box(slide, pt["detail"],
                      Inches(0.65), y + Inches(0.42),
                      W - Inches(1.1), Inches(0.6),
                      13, color=RGBColor(0xCC, 0xCC, 0xFF))

    # ── 5. Solution Slides ───────────────────────────────────────────────────
    for sol in content.get("solution_slides", []):
        slide = add_slide()
        light_bg(slide)
        header_bar(slide, sol["title"])
        footer(slide)

        # Value prop banner
        add_rect(slide, Inches(0.4), Inches(1.05), W - Inches(0.8), Inches(0.55),
                 ACCENT_BLUE)
        _add_text_box(slide, sol["value_prop"],
                      Inches(0.55), Inches(1.08),
                      W - Inches(1.1), Inches(0.5),
                      13, bold=True, color=WHITE)

        # Capabilities
        _add_text_box(slide, "Key Capabilities",
                      Inches(0.4), Inches(1.8),
                      Inches(5.5), Inches(0.4),
                      14, bold=True, color=SNAPLOGIC_DARK)
        txBox = slide.shapes.add_textbox(Inches(0.4), Inches(2.25), Inches(5.5), Inches(3.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        first = True
        for cap in sol.get("capabilities", []):
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(4)
            run = p.add_run()
            run.text = f"✓  {cap}"
            run.font.size = Pt(14)
            run.font.color.rgb = SNAPLOGIC_DARK

        # Outcome box
        add_rect(slide, Inches(6.2), Inches(1.8), Inches(6.7), Inches(4.0), SNAPLOGIC_PURPLE)
        _add_text_box(slide, "Customer Outcome",
                      Inches(6.4), Inches(1.95),
                      Inches(6.3), Inches(0.4),
                      13, bold=True, color=WHITE)
        _add_text_box(slide, sol["customer_outcome"],
                      Inches(6.4), Inches(2.5),
                      Inches(6.3), Inches(3.0),
                      14, color=RGBColor(0xDD, 0xDD, 0xFF))

    # ── 6. ROI / Business Impact Slide ──────────────────────────────────────
    slide = add_slide()
    purple_bg(slide)
    d = content["roi_slide"]
    header_bar(slide, d["title"], bg=SNAPLOGIC_DARK)
    footer(slide)

    metrics = d.get("metrics", [])
    col_count = len(metrics)
    col_w = (W - Inches(0.8)) / col_count if col_count else W
    for i, m in enumerate(metrics):
        x = Inches(0.4) + i * col_w
        add_rect(slide, x + Inches(0.1), Inches(1.1),
                 col_w - Inches(0.2), Inches(2.8), SNAPLOGIC_DARK)
        _add_text_box(slide, m["stat"],
                      x + Inches(0.2), Inches(1.3),
                      col_w - Inches(0.4), Inches(1.0),
                      32, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
        _add_text_box(slide, m["description"],
                      x + Inches(0.2), Inches(2.4),
                      col_w - Inches(0.4), Inches(1.2),
                      13, color=WHITE, align=PP_ALIGN.CENTER)

    _add_text_box(slide, d.get("summary", ""),
                  Inches(0.4), Inches(4.2),
                  W - Inches(0.8), Inches(1.2),
                  14, color=RGBColor(0xCC, 0xCC, 0xFF), align=PP_ALIGN.CENTER)

    # ── 7. Next Steps Slide ──────────────────────────────────────────────────
    slide = add_slide()
    light_bg(slide)
    d = content["next_steps_slide"]
    header_bar(slide, d["title"])
    footer(slide)

    steps = d.get("steps", [])
    step_h = (H - Inches(2.2)) / max(len(steps), 1)
    for i, s in enumerate(steps):
        y = Inches(1.1) + i * step_h
        # Number circle
        add_rect(slide, Inches(0.4), y + Inches(0.1),
                 Inches(0.55), Inches(0.55), SNAPLOGIC_PURPLE)
        _add_text_box(slide, str(i + 1),
                      Inches(0.4), y + Inches(0.05),
                      Inches(0.55), Inches(0.55),
                      16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Action
        _add_text_box(slide, s["action"],
                      Inches(1.1), y + Inches(0.05),
                      Inches(8.5), Inches(0.45),
                      15, bold=True, color=SNAPLOGIC_DARK)
        # Timeline badge
        add_rect(slide, Inches(10.0), y + Inches(0.1),
                 Inches(2.9), Inches(0.4), ACCENT_BLUE)
        _add_text_box(slide, s["timeline"],
                      Inches(10.0), y + Inches(0.08),
                      Inches(2.9), Inches(0.4),
                      12, color=WHITE, align=PP_ALIGN.CENTER)

    # ── 8. Closing Slide ─────────────────────────────────────────────────────
    slide = add_slide()
    purple_bg(slide)
    add_rect(slide, 0, 0, Inches(0.18), H, ACCENT_BLUE)
    d = content["closing_slide"]
    _add_text_box(slide, "SNAPLOGIC",
                  Inches(0.4), Inches(0.3),
                  Inches(4), Inches(0.5),
                  16, bold=True, color=WHITE)
    _add_text_box(slide, d["title"],
                  Inches(0.4), Inches(1.6),
                  W - Inches(1.0), Inches(1.4),
                  36, bold=True, color=WHITE)
    _add_text_box(slide, d["message"],
                  Inches(0.4), Inches(3.2),
                  W - Inches(1.0), Inches(1.4),
                  17, color=RGBColor(0xCC, 0xCC, 0xFF))
    _add_text_box(slide, d["contact_prompt"],
                  Inches(0.4), Inches(4.8),
                  W - Inches(1.0), Inches(0.5),
                  14, bold=True, color=ACCENT_BLUE)
    footer(slide, "Confidential — Prepared for " + customer_name)

    # ── Serialize ─────────────────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# -----------------------------
# UI
# -----------------------------
st.title("📊 Deck Builder")
st.markdown("Generate a customer-ready SnapLogic PowerPoint presentation in seconds.")
st.divider()

col_form, col_preview = st.columns([1, 1], gap="large")

with col_form:
    st.subheader("Customer Details")
    customer_name = st.text_input(
        "Customer / Prospect Name *",
        placeholder="e.g. Acme Corporation",
    )
    industry = st.selectbox("Industry *", INDUSTRIES)
    pain_points = st.text_area(
        "Key Pain Points / Challenges *",
        placeholder=(
            "e.g.\n"
            "- Manual, error-prone data processes across 12 systems\n"
            "- No real-time visibility into inventory levels\n"
            "- IT backlog is 6+ months for any integration work"
        ),
        height=160,
    )
    selected_solutions = st.multiselect(
        "Relevant SnapLogic Solutions *",
        SNAPLOGIC_SOLUTIONS,
        placeholder="Select one or more solutions...",
    )
    additional_context = st.text_area(
        "Additional Context (optional)",
        placeholder=(
            "e.g. They're currently using MuleSoft and looking to switch. "
            "Meeting is with the CTO and VP of Engineering."
        ),
        height=100,
    )

    st.divider()
    generate_btn = st.button(
        "Generate Presentation",
        type="primary",
        use_container_width=True,
        disabled=not (customer_name and pain_points and selected_solutions),
    )

with col_preview:
    st.subheader("What You'll Get")
    st.markdown("""
A branded, customer-ready `.pptx` with the following slides:

| # | Slide |
|---|-------|
| 1 | **Title** — Customer-specific cover |
| 2 | **Agenda** — Meeting roadmap |
| 3 | **Customer Challenges** — Their pain points framed professionally |
| 4 | **Why SnapLogic** — Platform overview & differentiation |
| 5+ | **Solution Deep-Dives** — One slide per selected solution |
| N-2 | **Business Impact** — ROI metrics tailored to the industry |
| N-1 | **Next Steps** — Actionable follow-up with timelines |
| N | **Closing** — Call to action |

All content is generated by Claude and tailored to the customer's industry and challenges.
""")
    st.info(
        "Tip: Be specific with pain points — the more detail you provide, "
        "the more tailored the deck.",
        icon="💡",
    )

    if "last_generated" in st.session_state:
        st.divider()
        st.success("Presentation ready! Download below.", icon="✅")
        st.download_button(
            label="Download PowerPoint (.pptx)",
            data=st.session_state["last_generated"]["bytes"],
            file_name=st.session_state["last_generated"]["filename"],
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
            type="primary",
        )

# -----------------------------
# Generation flow
# -----------------------------
if generate_btn:
    if not os.getenv("ANTHROPIC_API_KEY"):
        st.error(
            "ANTHROPIC_API_KEY is not set. Add it to your .env file or environment variables.",
            icon="🔑",
        )
        st.stop()

    with st.status("Building your presentation…", expanded=True) as status:
        try:
            st.write("Generating slide content with Claude...")
            slide_content = generate_slide_content(
                customer_name=customer_name,
                industry=industry,
                pain_points=pain_points,
                solutions=selected_solutions,
                additional_context=additional_context,
            )

            st.write("Building PowerPoint file...")
            pptx_bytes = build_presentation(slide_content, customer_name)

            safe_name = "".join(c if c.isalnum() or c in (" ", "-", "_") else "" for c in customer_name).strip()
            filename = f"SnapLogic_{safe_name}_Presentation.pptx"

            st.session_state["last_generated"] = {
                "bytes": pptx_bytes,
                "filename": filename,
            }

            status.update(label="Presentation ready!", state="complete")

        except json.JSONDecodeError as e:
            status.update(label="Failed", state="error")
            st.error(f"Claude returned unexpected content. Please try again. ({e})", icon="❌")
            st.stop()
        except anthropic.AuthenticationError:
            status.update(label="Failed", state="error")
            st.error("Invalid ANTHROPIC_API_KEY. Check your .env file.", icon="🔑")
            st.stop()
        except Exception as e:
            status.update(label="Failed", state="error")
            st.error(f"Something went wrong: {e}", icon="❌")
            st.stop()

    st.rerun()
